"""PowerPoint (.pptx) parser using python-pptx."""

from __future__ import annotations
import logging

logger = logging.getLogger(__name__)


async def parse_pptx(source: str) -> str:
    """Extract text from all slides of a .pptx file.

    Extracts slide titles, body text, and table content in order.
    """
    try:
        from pptx import Presentation
    except ImportError:
        logger.error("python-pptx not installed. Run: pip install python-pptx")
        return ""

    try:
        prs = Presentation(source)
        parts: list[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_parts: list[str] = [f"### Slide {slide_num}"]

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs).strip()
                        if text:
                            slide_parts.append(text)

                elif shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        if any(cells):
                            slide_parts.append(" | ".join(cells))

            if len(slide_parts) > 1:  # has content beyond the header
                parts.append("\n".join(slide_parts))

        return "\n\n".join(parts)

    except Exception as exc:
        logger.error("Failed to parse pptx '%s': %s", source, exc)
        return ""
